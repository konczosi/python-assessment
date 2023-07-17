import json
import logging
import pathlib
import sys
from json.decoder import JSONDecodeError

import numpy as np
from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from Task1_PPTX_report import pptx_utils

FILE_PATH_PREFIX = "Task1_PPTX_report/"


def generate_report(path: str) -> None:
    """Generate the presentation by the dict and save as output.pptx"""
    prs_config = read_json(path)
    try:
        sl_config = prs_config["presentation"]
        file_path = pathlib.Path("output.pptx")
        prs = Presentation()
        for sl in sl_config:
            generate_slide(sl, prs)
        prs.save(file_path)
        logging.info("Successfully generated %s report", file_path)
    except KeyError as error:
        logging.error("%s Key not found in %s", error, path)


def generate_slide(sl: dict, prs: Presentation) -> None:  # type: ignore
    """Generate the corresponding slide by type"""
    if sl["type"] == "title":
        generate_title_slide_content(sl, prs)
    elif sl["type"] == "text":
        generate_text_slide_content(sl, prs)
    elif sl["type"] == "list":
        generate_list_slide_content(sl, prs)
    elif sl["type"] == "picture":
        generate_picture_slide_content(sl, prs)
    elif sl["type"] == "plot":
        generate_chart_slide_content(sl, prs)
    else:
        logging.warning("There is no %s type slide", sl["type"])


def generate_title_slide_content(sl: dict, prs: Presentation) -> None:  # type: ignore
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = sl["title"]
    subtitle.text = sl["content"]
    logging.info("Successfully generated a %s type slide", sl["type"])


def generate_text_slide_content(sl: dict, prs: Presentation) -> None:  # type: ignore
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    title = slide.shapes.title
    title.text = sl["title"]
    left = Inches(0.5)
    top = Inches(1.75)
    width = Inches(9)
    height = Inches(4.95)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = sl["content"]
    logging.info("Successfully generated a %s type slide", sl["type"])


def generate_list_slide_content(sl: dict, prs: Presentation) -> None:  # type: ignore
    list_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(list_slide_layout)
    title = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title.text = sl["title"]
    tf = body_shape.text_frame
    i = 0
    for e in sl["content"]:
        tf.paragraphs[i].text = e["text"]
        tf.paragraphs[i].level = e["level"]
        tf.add_paragraph()
        i += 1
    logging.info("Successfully generated a %s type slide", sl["type"])


def generate_picture_slide_content(sl: dict, prs: Presentation) -> None:  # type: ignore
    img_path = FILE_PATH_PREFIX + sl["content"]
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    title = slide.shapes.title
    title.text = sl["title"]
    left = Inches(1.35)
    top = Inches(1.21)
    width = Inches(7.23)
    try:
        slide.shapes.add_picture(img_path, left, top, width=width)
        logging.info("Successfully generated a %s type slide", sl["type"])
    except FileNotFoundError as error:
        logging.error("Cannot add %s image to slide due to: %s", img_path, error)


def generate_chart_slide_content(sl: dict, prs: Presentation) -> None:  # type: ignore
    input_file_path = sl["content"].split(".")[0] + ".csv"
    pptx_utils.convert_csv_dat_binary(input_file_path)
    try:
        points = np.fromfile(sl["content"]).reshape(4, 2)
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        title.text = sl["title"]

        # Draw plot
        chart_data = XyChartData()
        series_1 = chart_data.add_series("")
        for p in points:
            series_1.add_data_point(p[0], p[1])
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS, x, y, cx, cy, chart_data  # type: ignore
        ).chart

        # Modify chart attributes
        chart.has_legend = False
        chart.value_axis.has_major_gridlines = False
        chart.category_axis.axis_title.text_frame.text = sl["configuration"]["x-label"]
        chart.value_axis.axis_title.text_frame.text = sl["configuration"]["y-label"]
        logging.info("Successfully generated a %s type slide", sl["type"])
    except (FileNotFoundError,ValueError) as error:
        logging.error("Opening %s file failed due to: %s", input_file_path, error)


def read_json(path: str) -> dict:
    """Returns a dictionary loaded from the json file"""
    config_file_path = pathlib.Path(path)
    try:
        with config_file_path.open("r") as file:
            config = json.load(file)
        return config
    except (FileNotFoundError, JSONDecodeError) as error:
        logging.error("Opening %s file failed due to: %s", config_file_path, error)
        sys.exit(1)
